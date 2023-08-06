import re
from os import PathLike
from pathlib import Path
from typing import Any, Dict, Optional, Union
from warnings import warn as warning

from PIL import Image
from pptx import Presentation

from .exceptions import RenderError
from .utils import fix_quotes, para_text_replace


class PPTXRenderer:
    """PPTX Renderer class

    This class is used to render a PPTX template by replacing python statements
    with the result of evaluating the python statements.

    Attributes:
        template_path (str): Path to the PPTX template.
    """

    def __init__(self, template_path: Union[str, bytes, PathLike]):
        self.template_path = template_path
        self.namespace = {}

    def render(
        self,
        output_path: Union[str, bytes, PathLike],
        methods_and_params: Optional[Dict[str, Any]] = None,
        skip_failed: bool = False,
    ) -> None:
        """Render PPTXRenderer template and save to output_path.

        Args:
            output_path (str): Path to the output PPTX file.
            methods_and_params (dict, optional): Dictionary of methods and parameters
                to be used in the template. Defaults to None.
            skip_failed (bool, optional): Dont raise an error if some of the
                statements failed to render. Defaults to False.

        Returns:
            None
        """
        if not Path(self.template_path).exists():
            raise (FileNotFoundError(f"{self.template_path} not found"))
        outppt = Presentation(self.template_path)
        self.namespace.update(methods_and_params)
        for slide_no, slide in enumerate(outppt.slides):
            if slide.has_notes_slide:
                python_code = re.search(
                    r"```python([\s\S]*)```",
                    fix_quotes(slide.notes_slide.notes_text_frame.text),
                    re.MULTILINE,
                )
                if python_code:
                    exec(python_code.group(1), self.namespace)
            for shape in list(slide.shapes):
                if shape.has_text_frame:
                    matches = re.finditer(r"{{{(.*)}}}", shape.text)
                    if not matches:
                        continue
                    for match_assignment in matches:
                        parts = match_assignment.group(1).split(":")
                        try:
                            result = eval(fix_quotes(parts[0]), self.namespace)
                        except Exception as ex:
                            if skip_failed:
                                warning(
                                    f"Evaluation of '{parts[0]}' in slide {slide_no+1} failed"
                                )
                                continue
                            raise RenderError(
                                f"Failed to evaluate '{parts[0]}'."
                            ) from ex
                        if len(parts) > 1 and parts[1].strip().lower() == "image":
                            if not Path(result).exists():
                                if skip_failed:
                                    warning(
                                        f"Image '{result}' in slide {slide_no+1} not found"
                                    )
                                    continue
                                raise RenderError(f"Image '{result}' not found.")
                            with Image.open(result) as img:
                                im_width, im_height = img.size
                            # add picture preserving aspect ratio
                            ar_image = im_width / im_height
                            ar_shape = shape.width / shape.height
                            if ar_image >= ar_shape:
                                slide.shapes.add_picture(
                                    result,
                                    shape.left,
                                    shape.top,
                                    shape.width,
                                    shape.width / ar_image,
                                )
                            else:
                                slide.shapes.add_picture(
                                    result,
                                    shape.left,
                                    shape.top,
                                    shape.height * ar_image,
                                    shape.height,
                                )
                            # Delete the shape after image is inserted
                            sp = shape._sp
                            sp.getparent().remove(sp)
                        elif len(parts) > 1 and parts[1].strip().lower() == "table":
                            try:
                                all_rows = list(result)
                                first_row_list = list(all_rows[0])
                                table_shape = slide.shapes.add_table(
                                    len(all_rows),
                                    len(first_row_list),
                                    shape.left,
                                    shape.top,
                                    shape.width,
                                    shape.height,
                                )
                                for row, row_data in enumerate(result):
                                    for col, val in enumerate(row_data):
                                        table_shape.table.cell(row, col).text = str(val)
                                # Delete the shape after image is inserted
                                sp = shape._sp
                                sp.getparent().remove(sp)
                            except Exception as ex:
                                if skip_failed:
                                    warning(
                                        f"Failed to render table {parts[0]} in slide {slide_no+1}"
                                    )
                                    continue
                                raise RenderError(
                                    f"Failed to render table from {parts[0]}."
                                ) from ex
                        else:
                            for paragraph in shape.text_frame.paragraphs:
                                para_text_replace(
                                    paragraph, match_assignment.group(0), result
                                )
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            matches = re.finditer(r"{{{(.*)}}}", cell.text)
                            if not matches:
                                continue
                            for match_assignment in matches:
                                parts = match_assignment.group(1).split(":")
                                try:
                                    result = eval(fix_quotes(parts[0]), self.namespace)
                                except Exception as ex:
                                    if skip_failed:
                                        warning(
                                            f"Evaluation of '{parts[0]}' in slide {slide_no+1} failed"
                                        )
                                        continue
                                    raise RenderError(
                                        f"Failed to evaluate '{parts[0]}'."
                                    ) from ex
                                for paragraph in cell.text_frame.paragraphs:
                                    para_text_replace(
                                        paragraph, match_assignment.group(0), result
                                    )
        outppt.save(output_path)
